#!/usr/bin/env python3
"""
Vault Integration Architecture - PowerPoint Generator
Run on Windows: python generate_vault_ppt.py
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_vault_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Vault Integration Architecture"
    subtitle.text = "demoApp - Spring Boot & HashiCorp Vault kv-v2"
    
    # Slide 2: Current Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Current Architecture (As-Is - 2026-05-07)"
    content = slide.placeholders[1]
    content.text = """Vault Server: 192.168.2.57:8200
Mount: ebiz_service (kv-v2)
Secret Path: ebiz_db/data-enc-key
Key: fernet-key (32 bytes)

Flow:
1. demoApp → Spring Cloud Vault → HashiCorp Vault
2. PasswordService reads Fernet key from Vault
3. AES-128 encryption with CBC mode
4. Base64 encoded output stored in database"""
    
    # Slide 3: Vault kv-2 Path Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Vault kv-2 Secret Structure"
    content = slide.placeholders[1]
    content.text = """Vault CLI Commands:

# Enable kv-v2 backend
vault secrets enable -path=ebiz_service kv-v2

# Store Fernet key
vault kv put -mount=ebiz_service ebiz_db/data-enc-key \\
  fernet-key="NgqOBievnB9500cQOnSQ-..." \\
  description="encryption key for ebiz db column"

# Read secret (kv-v2 adds 'data/' automatically)
vault kv get -mount=ebiz_service ebiz_db/data-enc-key

API Path: /v1/ebiz_service/data/ebiz_db/data-enc-key"""
    
    # Slide 4: Fernet Key Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Fernet Key Structure (32 bytes total)"
    content = slide.placeholders[1]
    content.text = """Bytes 0-15: AES-128 Key (16 bytes)
Bytes 16-31: HMAC-SHA256 Key (16 bytes)

Format: URL-safe Base64 (uses '-' and '_')
Example: NgqOBievnB9500cQOnSQ-cmbBx38KnOiKx5ooQ_e97Y=

Decoded: 32 bytes binary data
- First 16 bytes → AES-128 secret key
- Last 16 bytes → HMAC-SHA256 key (currently unused)"""
    
    # Slide 5: PasswordService Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "VaultOperations.read() Integration"
    content = slide.placeholders[1]
    content.text = """@Service
public class PasswordService {
    private byte[] encryptionKey;
    
    public PasswordService(VaultOperations vaultOps) {
        VaultResponse response = vaultOps.read(
            "ebiz_service/data/ebiz_db/data-enc-key");
        
        Map<String, Object> outerData = response.getData();
        Map<String, Object> secretData = 
            (Map<String, Object>) outerData.get("data");
        
        String fernetKeyBase64 = 
            (String) secretData.get("fernet-key");
        
        // Decode URL-safe Base64
        this.encryptionKey = 
            Base64.getUrlDecoder().decode(fernetKeyBase64);
    }
}"""
    
    # Slide 6: Encryption Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Password Encryption Flow"
    content = slide.placeholders[1]
    content.text = """User Input: "vaulttest123"
     ↓
PasswordService.encryptPassword()
     ↓
Fernet Key (byte[32]) → AES-128 (first 16 bytes)
     ↓
Cipher.doFinal(password.getBytes())
     ↓
Base64.getEncoder().encodeToString()
     ↓
Output: "pU4nAaBrwqPKLoV1Waa/tw=="
     ↓
Stored in DB: ebiz.board.password"""
    
    # Slide 7: VaultResponse Structure (kv-v2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Vault Read Response (kv-v2 format)"
    content = slide.placeholders[1]
    content.text = """{
  "data": {
    "data": {
      "description": "encryption key for ebiz db column",
      "fernet-key": "NgqOBievnB9500cQOnSQ-..."
    },
    "metadata": {
      "created_time": "2026-05-06T11:32:48...",
      "version": 1
    }
  }
}

Key Points:
- kv-v2 wraps data in double "data" fields
- Outer: response.getData()
- Inner: outerData.get("data")"""
    
    # Slide 8: Configuration Files
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "application.properties Configuration"
    content = slide.placeholders[1]
    content.text = """# Vault Configuration
spring.cloud.vault.uri=http://192.168.2.57:8200
spring.cloud.vault.token=${VAULT_TOKEN}
spring.cloud.vault.fail-fast=false

Key Settings:
- fail-fast=false: App starts even if Vault is down
- Token: Stored in app.properties (file is in .gitignore)
- URI: Vault server address

Note: No spring.cloud.vault.kv.* properties needed 
      since we use VaultOperations.read() directly"""
    
    # Slide 9: Verification Results
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Verification (Tested 2026-05-07)"
    content = slide.placeholders[1]
    content.text = """✅ Test Results:
- App Startup: SUCCESS (reads Fernet key from Vault)
- POST /api/v1/posts: SUCCESS (post ID 2017587 created)
- Password Encryption: SUCCESS 
  → "pU4nAaBrwqPKLoV1Waa/tw==" stored in DB
- Vault Unavailable: Graceful (fail-fast=false)

Database Verification:
SELECT id, title, password FROM ebiz.board 
WHERE id = 2017587;
→ password = "pU4nAaBrwqPKLoV1Waa/tw==" (encrypted)"""
    
    # Slide 10: Security Considerations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Security Best Practices"
    content = slide.placeholders[1]
    content.text = """✅ Implemented:
1. Fernet key read from Vault kv-v2 at startup
2. URL-safe Base64 decoding for Fernet format
3. AES-128 encryption with CBC mode
4. Base64 output for database storage

⚠️ To Improve:
1. Move Vault token to environment variable
2. Enable TLS/HTTPS for Vault communication (currently HTTP)
3. Use AppRole authentication instead of Token for production
4. Implement Fernet key rotation strategy
5. Enable Vault audit logging
6. Restrict network access to Vault (firewall rules)"""
    
    # Slide 11: Files Modified
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Repository Changes"
    content = slide.placeholders[1]
    content.text = """Files Modified:
1. PasswordService.java - Vault-based encryption key loading
2. application.properties - Vault connection settings
3. VAULT_INTEGRATION_DIAGRAM.md - Architecture diagrams
4. VAULT_AND_ENCRYPTION.md - Documentation

GitHub: https://github.com/yadalhan/demoApp

Commit: 7c1b2cb - Update Vault integration diagram
Commit: d304543 - Fix Vault kv-v2 integration with Fernet key"""
    
    # Slide 12: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Questions?"
    subtitle.text = "Thank You!\nSisyphus - Powered by OhMyOpenCode"
    
    # Save presentation
    output_file = "Vault_Integration_Architecture.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    try:
        create_vault_presentation()
    except ImportError:
        print("❌ Error: python-pptx not installed")
        print("   Install: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error: {e}")
